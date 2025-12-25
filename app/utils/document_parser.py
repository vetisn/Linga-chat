"""
文档解析工具 - 支持多种文件格式
支持: PDF, DOCX, DOC, PPTX, XLSX, TXT, MD, CSV
"""
import os
from typing import Optional


def extract_text_from_file(file_path: str) -> str:
    """
    根据文件扩展名自动选择解析方法提取文本
    """
    ext = os.path.splitext(file_path)[1].lower()
    
    extractors = {
        '.pdf': extract_pdf,
        '.docx': extract_docx,
        '.doc': extract_doc,
        '.pptx': extract_pptx,
        '.xlsx': extract_xlsx,
        '.xls': extract_xlsx,
        '.txt': extract_text,
        '.md': extract_text,
        '.csv': extract_text,
        '.json': extract_text,
        '.xml': extract_text,
        '.html': extract_html,
        '.htm': extract_html,
    }
    
    extractor = extractors.get(ext)
    if not extractor:
        raise ValueError(f"不支持的文件格式: {ext}")
    
    return extractor(file_path)


def extract_pdf(file_path: str) -> str:
    """提取 PDF 文本，支持扫描件OCR"""
    try:
        from PyPDF2 import PdfReader
        reader = PdfReader(file_path)
        text_parts = []
        has_text = False
        
        for page in reader.pages:
            page_text = page.extract_text()
            if page_text and page_text.strip():
                text_parts.append(page_text)
                has_text = True
        
        # 如果提取到足够的文本，直接返回
        if has_text and len("".join(text_parts)) > 100:
            return "\n\n".join(text_parts)
        
        # 尝试使用OCR处理扫描件PDF
        ocr_text = extract_pdf_with_ocr(file_path)
        if ocr_text:
            return ocr_text
        
        # 如果OCR也没有结果，返回已提取的文本
        if text_parts:
            return "\n\n".join(text_parts)
        
        raise ValueError("PDF 中未找到可提取的文本，可能是扫描件。请安装 Tesseract OCR 以支持扫描件识别。")
        
    except ImportError:
        raise ImportError("请安装 PyPDF2: pip install PyPDF2")
    except Exception as e:
        raise ValueError(f"PDF 解析失败: {e}")


def extract_pdf_with_ocr(file_path: str) -> Optional[str]:
    """使用 Tesseract OCR 提取扫描件 PDF 中的文字"""
    try:
        import pytesseract
        from pdf2image import convert_from_path
        
        # 检查 Tesseract 是否可用
        try:
            pytesseract.get_tesseract_version()
        except Exception:
            return None
        
        # 将 PDF 转换为图片
        images = convert_from_path(file_path, dpi=200)
        
        text_parts = []
        for i, image in enumerate(images):
            # 使用 Tesseract 识别文字，支持中英文
            page_text = pytesseract.image_to_string(image, lang='chi_sim+eng')
            if page_text and page_text.strip():
                text_parts.append(f"[第 {i+1} 页]\n{page_text.strip()}")
        
        if text_parts:
            return "\n\n".join(text_parts)
        return None
        
    except ImportError:
        # pytesseract 或 pdf2image 未安装
        return None
    except Exception as e:
        print(f"OCR 处理失败: {e}")
        return None


def extract_docx(file_path: str) -> str:
    """提取 DOCX 文本"""
    try:
        from docx import Document
        doc = Document(file_path)
        text_parts = []
        for para in doc.paragraphs:
            if para.text.strip():
                text_parts.append(para.text)
        # 也提取表格内容
        for table in doc.tables:
            for row in table.rows:
                row_text = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                if row_text:
                    text_parts.append(" | ".join(row_text))
        return "\n".join(text_parts)
    except ImportError:
        raise ImportError("请安装 python-docx: pip install python-docx")
    except Exception as e:
        raise ValueError(f"DOCX 解析失败: {e}")


def extract_doc(file_path: str) -> str:
    """提取 DOC 文本 (旧版 Word 格式)"""
    # DOC 格式比较复杂，尝试几种方法
    # 方法1: 尝试用 docx 打开（有时可以）
    try:
        return extract_docx(file_path)
    except:
        pass
    
    # 方法2: 尝试用 antiword (如果系统安装了)
    try:
        import subprocess
        result = subprocess.run(['antiword', file_path], capture_output=True, text=True)
        if result.returncode == 0:
            return result.stdout
    except:
        pass
    
    # 方法3: 尝试读取为文本（某些 DOC 文件可能包含可读文本）
    try:
        with open(file_path, 'rb') as f:
            content = f.read()
            # 尝试提取可打印字符
            import re
            text = content.decode('utf-8', errors='ignore')
            # 过滤掉大部分乱码
            text = re.sub(r'[^\x20-\x7E\u4e00-\u9fff\n]', ' ', text)
            text = re.sub(r'\s+', ' ', text).strip()
            if len(text) > 100:
                return text
    except:
        pass
    
    raise ValueError("DOC 格式解析失败，建议转换为 DOCX 格式后上传")


def extract_pptx(file_path: str) -> str:
    """提取 PPTX 文本"""
    try:
        from pptx import Presentation
        prs = Presentation(file_path)
        text_parts = []
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_texts = [f"[幻灯片 {slide_num}]"]
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_texts.append(shape.text)
                # 提取表格
                if shape.has_table:
                    for row in shape.table.rows:
                        row_text = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                        if row_text:
                            slide_texts.append(" | ".join(row_text))
            if len(slide_texts) > 1:
                text_parts.append("\n".join(slide_texts))
        return "\n\n".join(text_parts)
    except ImportError:
        raise ImportError("请安装 python-pptx: pip install python-pptx")
    except Exception as e:
        raise ValueError(f"PPTX 解析失败: {e}")


def extract_xlsx(file_path: str) -> str:
    """提取 XLSX/XLS 文本"""
    try:
        from openpyxl import load_workbook
        wb = load_workbook(file_path, data_only=True)
        text_parts = []
        for sheet_name in wb.sheetnames:
            sheet = wb[sheet_name]
            sheet_texts = [f"[工作表: {sheet_name}]"]
            for row in sheet.iter_rows():
                row_values = []
                for cell in row:
                    if cell.value is not None:
                        row_values.append(str(cell.value))
                if row_values:
                    sheet_texts.append(" | ".join(row_values))
            if len(sheet_texts) > 1:
                text_parts.append("\n".join(sheet_texts))
        return "\n\n".join(text_parts)
    except ImportError:
        raise ImportError("请安装 openpyxl: pip install openpyxl")
    except Exception as e:
        raise ValueError(f"XLSX 解析失败: {e}")


def extract_text(file_path: str) -> str:
    """提取纯文本文件"""
    # 尝试检测编码
    try:
        import chardet
        with open(file_path, 'rb') as f:
            raw = f.read()
            detected = chardet.detect(raw)
            encoding = detected.get('encoding', 'utf-8') or 'utf-8'
        return raw.decode(encoding)
    except ImportError:
        pass
    
    # 回退到常见编码尝试
    encodings = ['utf-8', 'gbk', 'gb2312', 'utf-16', 'latin-1']
    for enc in encodings:
        try:
            with open(file_path, 'r', encoding=enc) as f:
                return f.read()
        except UnicodeDecodeError:
            continue
    
    raise ValueError("无法识别文件编码")


def extract_html(file_path: str) -> str:
    """提取 HTML 文本"""
    import re
    content = extract_text(file_path)
    # 移除 script 和 style 标签
    content = re.sub(r'<script[^>]*>.*?</script>', '', content, flags=re.DOTALL | re.IGNORECASE)
    content = re.sub(r'<style[^>]*>.*?</style>', '', content, flags=re.DOTALL | re.IGNORECASE)
    # 移除所有 HTML 标签
    content = re.sub(r'<[^>]+>', ' ', content)
    # 清理多余空白
    content = re.sub(r'\s+', ' ', content).strip()
    return content


def get_supported_extensions() -> list:
    """返回支持的文件扩展名列表"""
    return ['.pdf', '.docx', '.doc', '.pptx', '.xlsx', '.xls', '.txt', '.md', '.csv', '.json', '.xml', '.html', '.htm']
